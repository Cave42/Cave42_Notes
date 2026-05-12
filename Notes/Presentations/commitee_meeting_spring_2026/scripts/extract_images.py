#!/usr/bin/env python3
from pptx import Presentation
import os

PPTX_PATH = os.path.expanduser('~/Downloads/committee_meeting_spring.pptx')
PUBLIC_DIR = os.path.join(os.path.dirname(__file__), '..', 'public')

CONTENT_TYPE_EXT = {
    'image/png': '.png',
    'image/jpeg': '.jpg',
    'image/gif': '.gif',
    'image/bmp': '.bmp',
}

prs = Presentation(PPTX_PATH)

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    images = [s for s in slide.shapes if s.shape_type == 13]
    for j, shape in enumerate(images):
        try:
            image = shape.image
        except ValueError:
            continue
        ext = CONTENT_TYPE_EXT.get(image.content_type, '.png')
        filename = f'slide_{slide_num}{ext}' if len(images) == 1 else f'slide_{slide_num}_{j+1}{ext}'
        filepath = os.path.join(PUBLIC_DIR, filename)
        with open(filepath, 'wb') as f:
            f.write(image.blob)
        print(f'Saved {filename}')

print('Done.')
